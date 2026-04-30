"""Generate final presentation PowerPoint — 15 slides, 15 minutes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MED = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_RED = RGBColor(0xF4, 0x43, 0x36)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

FIG_DIR = os.path.join(os.path.dirname(__file__), "figures")


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_stat_box(slide, left, top, width, height, number, label,
                 num_color=ACCENT_BLUE, bg_color=None):
    if bg_color:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

    add_text_box(slide, left, top + Inches(0.1), width, Inches(0.6),
                 number, font_size=36, color=num_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.7), width, Inches(0.4),
                 label, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 1: Title
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl)

add_text_box(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "An Empirical Evaluation of Vision-Language Model\nVisualization Literacy Across 18 Chart Types",
             font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(3.3), Inches(11), Inches(0.5),
             "Insights for Immersive Analytics",
             font_size=22, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
             "Naga Venkata Sai Chennu  |  Hemanjali Buchireddy",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
             "CS 692: Mobile Immersive Computing  |  George Mason University",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Advisor: Dr. Bo Han  |  Collaborator: Fahim Arsad Nafis",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             "April 17, 2026",
             font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 2: Project Overview (1 recap slide)
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Project Overview", font_size=28, color=WHITE, bold=True)

# Left: Research question
add_text_box(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Research Question", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(1),
             "How well can Vision-Language Models understand data visualizations \u2014 "
             "and what happens when charts move from 2D to 3D (immersive)?",
             font_size=16, color=WHITE)

add_text_box(sl, Inches(0.8), Inches(3), Inches(5.5), Inches(0.5),
             "What We Did", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(sl, Inches(0.8), Inches(3.5), Inches(5.5), Inches(2.5), [
    "\u2022  Evaluated GPT-5.4 on ChartX benchmark",
    "\u2022  6,000 chart images \u00d7 18 chart types \u00d7 22 topics",
    "\u2022  5,720 2D images + 280 3D bar chart images",
    "\u2022  Original images & QA pairs from the dataset",
    "\u2022  Zero-shot, temperature=0, fully automated",
], font_size=14)

# Right: Key numbers
add_text_box(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Key Numbers", font_size=18, color=ACCENT_BLUE, bold=True)

add_stat_box(sl, Inches(7), Inches(2), Inches(2.5), Inches(1.1),
             "85.7%", "2D Accuracy", ACCENT_GREEN, RGBColor(0x20, 0x30, 0x50))
add_stat_box(sl, Inches(9.8), Inches(2), Inches(2.5), Inches(1.1),
             "59.3%", "3D Accuracy", ACCENT_RED, RGBColor(0x20, 0x30, 0x50))
add_stat_box(sl, Inches(7), Inches(3.5), Inches(2.5), Inches(1.1),
             "43.7pp", "Type Gap", ACCENT_ORANGE, RGBColor(0x20, 0x30, 0x50))
add_stat_box(sl, Inches(9.8), Inches(3.5), Inches(2.5), Inches(1.1),
             "6,000", "Images Tested", ACCENT_BLUE, RGBColor(0x20, 0x30, 0x50))

add_text_box(sl, Inches(7), Inches(5), Inches(5.5), Inches(0.5),
             "Why It Matters", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(sl, Inches(7), Inches(5.5), Inches(5.5), Inches(1.2),
             "Immersive analytics (VR/AR) uses 3D charts by default. "
             "If VLMs can\u2019t read 3D charts, they can\u2019t assist users in VR. "
             "No prior work evaluates VLMs across this many chart types.",
             font_size=14, color=LIGHT_GRAY)


# ==========================================================================
# SLIDE 3: ChartX Dataset
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "The ChartX Benchmark Dataset", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Dataset Details", font_size=18, color=ACCENT_BLUE, bold=True)

# Dataset stats
stats = [
    ("6,000", "chart images (PNG)"),
    ("18", "chart types"),
    ("22", "academic topics"),
    ("1", "QA pair per image (human-authored)"),
    ("17", "2D types + 1 3D type (3D-Bar)"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.9 + i * 0.5)
    add_text_box(sl, Inches(1), y, Inches(1.2), Inches(0.4),
                 num, font_size=20, color=ACCENT_BLUE, bold=True)
    add_text_box(sl, Inches(2.3), y, Inches(4), Inches(0.4),
                 label, font_size=16, color=WHITE)

add_text_box(sl, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5),
             "Sample Question", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(sl, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.4),
             "Q: How many more nurses are in East vs North?",
             font_size=14, color=WHITE, bold=True)
add_text_box(sl, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.4),
             "A: 500", font_size=14, color=ACCENT_GREEN, bold=True)

# Right: Sample image
if os.path.exists(f"{FIG_DIR}/sample_bar_2d.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/sample_bar_2d.png",
                          Inches(7), Inches(1.2), Inches(5.5), Inches(4))

add_text_box(sl, Inches(7), Inches(5.5), Inches(5.5), Inches(0.5),
             "Source: HuggingFace (InternScience/ChartX)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 4: Methodology
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Methodology", font_size=28, color=WHITE, bold=True)

# Pipeline steps
steps = [
    ("1", "Load", "Chart image + QA pair from ChartX"),
    ("2", "Query", "Send image + question to GPT-5.4 (temp=0)"),
    ("3", "Parse", "Extract answer from free-text response"),
    ("4", "Score", "10% numeric tolerance + keyword matching"),
    ("5", "Cache", "JSON response files for reproducibility"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.3 + i * 0.9)
    # Number circle
    shape = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(sl, Inches(1.8), y, Inches(1.5), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(sl, Inches(3.3), y + Inches(0.05), Inches(4), Inches(0.4),
                 desc, font_size=14, color=LIGHT_GRAY)

# Right side: config
add_text_box(sl, Inches(7.5), Inches(1.3), Inches(5), Inches(0.5),
             "Configuration", font_size=20, color=ACCENT_BLUE, bold=True)

config = [
    "Model: GPT-5.4 (OpenAI)",
    "Temperature: 0 (deterministic)",
    "Max tokens: 200",
    "Concurrency: 5 parallel API calls",
    "Retry: Exponential backoff",
    "Total cost: ~$10",
    "Total evaluations: 5,713 (2D) + 280 (3D)",
]
add_bullet_list(sl, Inches(7.5), Inches(2), Inches(5), Inches(3.5), config, font_size=14)

add_text_box(sl, Inches(7.5), Inches(5.5), Inches(5), Inches(1),
             "Zero-shot evaluation \u2014 no system prompt,\nno few-shot examples, no chain-of-thought",
             font_size=14, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 5: KEY RESULT \u2014 2D Accuracy Table
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Result: 2D Accuracy by Chart Type", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "Overall: 85.7% (4,896 / 5,713)  |  17 chart types  |  GPT-5.4",
             font_size=16, color=ACCENT_BLUE, bold=True)

if os.path.exists(f"{FIG_DIR}/fig1_accuracy_by_chart_type.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/fig1_accuracy_by_chart_type.png",
                          Inches(0.5), Inches(1.5), Inches(8.5), Inches(5.5))

# Key insight box
add_text_box(sl, Inches(9.3), Inches(1.8), Inches(3.5), Inches(0.5),
             "Key Insight", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(sl, Inches(9.3), Inches(2.5), Inches(3.5), Inches(4),
             "13 chart types > 83%\n\n"
             "4 chart types struggle:\n"
             "\u2022 Area chart: 77.4%\n"
             "\u2022 Bubble: 71.1%\n"
             "\u2022 Radar: 68.2%\n"
             "\u2022 Treemap: 51.6%\n\n"
             "43.7pp gap between\nbest and worst type",
             font_size=14, color=WHITE)


# ==========================================================================
# SLIDE 6: The 43-Point Gap
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "The Perceptual Encoding Hierarchy", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1), Inches(11), Inches(0.8),
             "VLM accuracy follows the same hierarchy as human vision (Cleveland & McGill, 1984)",
             font_size=18, color=LIGHT_GRAY)

if os.path.exists(f"{FIG_DIR}/fig3_perceptual_hierarchy.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/fig3_perceptual_hierarchy.png",
                          Inches(0.5), Inches(2), Inches(7), Inches(4.5))

add_text_box(sl, Inches(8), Inches(2), Inches(4.8), Inches(0.5),
             "Cleveland & McGill (1984)", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(sl, Inches(8), Inches(2.6), Inches(4.8), Inches(4),
             "Most accurate:\n"
             "  1. Position on common scale\n"
             "  2. Position on non-aligned scales\n"
             "  3. Length\n"
             "  4. Direction / Angle\n"
             "  5. Area\n"
             "  6. Volume\n"
             "  7. Color saturation\n\n"
             "Our results match this\nhierarchy almost perfectly.",
             font_size=14, color=WHITE)


# ==========================================================================
# SLIDE 7: 3D Collapse
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Result: The 3D Accuracy Collapse", font_size=28, color=WHITE, bold=True)

add_stat_box(sl, Inches(1), Inches(1.5), Inches(3), Inches(1.3),
             "92.1%", "2D bar_chart_num", ACCENT_GREEN, RGBColor(0x20, 0x30, 0x50))
add_stat_box(sl, Inches(4.5), Inches(1.5), Inches(3), Inches(1.3),
             "87.7%", "2D bar_chart", ACCENT_GREEN, RGBColor(0x20, 0x30, 0x50))
add_stat_box(sl, Inches(8), Inches(1.5), Inches(3), Inches(1.3),
             "59.3%", "3D-Bar", ACCENT_RED, RGBColor(0x20, 0x30, 0x50))

add_text_box(sl, Inches(1), Inches(3.2), Inches(10), Inches(0.5),
             "\u2193 26\u201333 percentage-point drop from 2D to 3D",
             font_size=20, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

# Sample images
if os.path.exists(f"{FIG_DIR}/sample_bar_2d.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/sample_bar_2d.png",
                          Inches(1.5), Inches(4), Inches(4.5), Inches(3))
if os.path.exists(f"{FIG_DIR}/sample_bar_3d.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/sample_bar_3d.png",
                          Inches(7), Inches(4), Inches(4.5), Inches(3))

add_text_box(sl, Inches(1.5), Inches(3.7), Inches(4.5), Inches(0.4),
             "2D Bar Chart", font_size=14, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(7), Inches(3.7), Inches(4.5), Inches(0.4),
             "3D Bar Chart", font_size=14, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 8: Why 3D Fails
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Why 3D Charts Fail \u2014 Four Distortion Mechanisms", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "Based on Wilke (2019) \u2014 Fundamentals of Data Visualization",
             font_size=14, color=LIGHT_GRAY)

mechanisms = [
    ("1. Non-invertible Projection",
     "A single 2D pixel maps to a line in 3D space.\nExact values cannot be recovered from one view."),
    ("2. Foreshortening",
     "Bars farther from the camera appear shorter\nthan equally-valued bars in front."),
    ("3. Occlusion",
     "Front elements hide rear elements,\nphysically removing information from the image."),
    ("4. Unequal Scaling",
     "Perspective causes identical data increments\nto occupy different pixel heights at different depths."),
]
for i, (title, desc) in enumerate(mechanisms):
    y = Inches(1.7 + i * 1.3)
    add_text_box(sl, Inches(1), y, Inches(5), Inches(0.4),
                 title, font_size=18, color=ACCENT_RED, bold=True)
    add_text_box(sl, Inches(1), y + Inches(0.4), Inches(5), Inches(0.7),
                 desc, font_size=14, color=WHITE)

add_text_box(sl, Inches(7), Inches(1.7), Inches(5.5), Inches(5),
             "VLMs see a single static 2D\nprojection of a 3D chart.\n\n"
             "They cannot:\n"
             "\u2022  Rotate the view\n"
             "\u2022  Look behind occluded bars\n"
             "\u2022  Judge depth from stereo vision\n"
             "\u2022  Interact with the chart\n\n"
             "Humans in VR can compensate.\nVLMs cannot.",
             font_size=16, color=LIGHT_GRAY)


# ==========================================================================
# SLIDE 9: Invisible Ground Truth
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Key Finding: The Invisible Ground Truth Problem", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
             "Why does treemap accuracy = 51.6%?", font_size=18, color=ACCENT_ORANGE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
             "The image shows:", font_size=16, color=WHITE, bold=True)
add_text_box(sl, Inches(0.8), Inches(2.1), Inches(5), Inches(0.5),
             "Colored rectangles + category names\nNO numeric values", font_size=14, color=WHITE)

add_text_box(sl, Inches(0.8), Inches(2.8), Inches(5), Inches(0.4),
             "The ground truth expects:", font_size=16, color=WHITE, bold=True)
add_text_box(sl, Inches(0.8), Inches(3.2), Inches(5), Inches(0.5),
             "Exact percentages (e.g., \"20%\")\nfrom CSV data NOT shown in image", font_size=14, color=ACCENT_RED)

add_text_box(sl, Inches(0.8), Inches(4), Inches(5), Inches(0.4),
             "The reality:", font_size=16, color=WHITE, bold=True)
add_text_box(sl, Inches(0.8), Inches(4.4), Inches(5), Inches(1.5),
             "Stevens' Power Law: area perception\n"
             "exponent a \u2248 0.7 (compressive)\n\n"
             "Even humans cannot extract exact\n"
             "values from area encoding without labels.\n\n"
             "This is a benchmark design issue,\nnot a VLM failure.",
             font_size=14, color=LIGHT_GRAY)

if os.path.exists(f"{FIG_DIR}/sample_treemap.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/sample_treemap.png",
                          Inches(6.5), Inches(1.2), Inches(6), Inches(4))

add_text_box(sl, Inches(6.5), Inches(5.5), Inches(6), Inches(0.8),
             "ChartX paper acknowledges: \"the column label of\n"
             "values is usually invisible\" for treemaps",
             font_size=12, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 10: Failure Analysis
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Failure Analysis: Why Weak Chart Types Fail", font_size=28, color=WHITE, bold=True)

if os.path.exists(f"{FIG_DIR}/fig4_failure_analysis.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/fig4_failure_analysis.png",
                          Inches(0.5), Inches(1.2), Inches(6), Inches(4.5))

add_text_box(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Root Causes", font_size=20, color=ACCENT_BLUE, bold=True)

causes = [
    ("Treemap (51.6%)", "No numbers visible in image.\nMust guess % from rectangle area."),
    ("Radar (68.2%)", "Overlapping series on polar axes.\n22% are format mismatches."),
    ("Bubble (71.1%)", "Size encoding is ambiguous.\nCan't read values from bubble diameter."),
    ("Area (77.4%)", "Stacked fills confuse individual\nvs cumulative values."),
]
for i, (ct, cause) in enumerate(causes):
    y = Inches(2 + i * 1.2)
    add_text_box(sl, Inches(7), y, Inches(5.5), Inches(0.4),
                 ct, font_size=14, color=ACCENT_RED, bold=True)
    add_text_box(sl, Inches(7), y + Inches(0.35), Inches(5.5), Inches(0.6),
                 cause, font_size=12, color=LIGHT_GRAY)


# ==========================================================================
# SLIDE 11: Design Guidelines
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Six Design Guidelines for VLM-Assisted Analytics", font_size=28, color=WHITE, bold=True)

guidelines = [
    ("G1", "Use position-encoded charts", "(bar, line) for VLM tasks \u2014 90%+ accuracy"),
    ("G2", "Always include numeric labels", "Transforms estimation \u2192 text reading (+4.4pp)"),
    ("G3", "Flatten 3D before VLM processing", "Render 2D projections or provide data tables"),
    ("G4", "Tiered scoring tolerance", "5% for bar/line, 15% for pie, 25% for treemap"),
    ("G5", "Describe \u2192 Extract \u2192 Reason", "Structure prompts as multi-step pipelines"),
    ("G6", "VLM-interpretable layer in VR", "Pass DXR JSON specs alongside viewport images"),
]
for i, (num, title, desc) in enumerate(guidelines):
    y = Inches(1.2 + i * 0.95)
    # Number
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.7), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(sl, Inches(1.7), y, Inches(5), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(sl, Inches(1.7), y + Inches(0.35), Inches(5), Inches(0.4),
                 desc, font_size=13, color=LIGHT_GRAY)


# ==========================================================================
# SLIDE 12: DXR & Immersive Analytics Connection
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Connection to Immersive Analytics (DXR Toolkit)", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "DXR: Declarative Immersive Visualization", font_size=18, color=ACCENT_BLUE, bold=True)

add_bullet_list(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3), [
    "\u2022  Unity toolkit for VR/AR data visualization",
    "\u2022  Vega-Lite-inspired JSON specification",
    "\u2022  20 mark types (cube, sphere, bar, etc.)",
    "\u2022  Deploys on Meta Quest headsets",
    "\u2022  Renders 3D bar, scatter, heatmap, etc.",
], font_size=14)

add_text_box(sl, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5),
             "How Our Results Apply", font_size=18, color=ACCENT_BLUE, bold=True)

add_bullet_list(sl, Inches(0.8), Inches(5.1), Inches(5.5), Inches(2), [
    "\u2022  DXR renders 3D charts \u2192 VLM sees viewport screenshot",
    "\u2022  Our 3D results (59.3%) predict VLM accuracy in VR",
    "\u2022  G6: Pass DXR JSON specs as structured context",
    "\u2022  G3: Flatten 3D \u2192 2D before querying VLM",
], font_size=14)

add_text_box(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Proposed VLM + DXR Architecture", font_size=18, color=ACCENT_BLUE, bold=True)

arch_steps = [
    "User sees 3D chart in VR (DXR)",
    "\u2193",
    "VLM receives: viewport image + DXR JSON",
    "\u2193",
    "VLM answers user's question",
    "\u2193",
    "Answer displayed in VR overlay",
]
for i, step in enumerate(arch_steps):
    y = Inches(1.9 + i * 0.55)
    color = ACCENT_BLUE if step == "\u2193" else WHITE
    sz = 20 if step == "\u2193" else 14
    add_text_box(sl, Inches(7.5), y, Inches(5), Inches(0.4),
                 step, font_size=sz, color=color, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 13: Comparison with Prior Work
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Comparison with Prior Work", font_size=28, color=WHITE, bold=True)

if os.path.exists(f"{FIG_DIR}/fig5_prior_work.png"):
    sl.shapes.add_picture(f"{FIG_DIR}/fig5_prior_work.png",
                          Inches(0.5), Inches(1.2), Inches(7), Inches(4.5))

add_text_box(sl, Inches(8), Inches(1.2), Inches(4.8), Inches(0.5),
             "Our Contributions", font_size=18, color=ACCENT_BLUE, bold=True)

add_bullet_list(sl, Inches(8), Inches(1.9), Inches(4.8), Inches(4.5), [
    "\u2022  Broadest chart type diversity\n   (18 types vs 6\u201312 in prior work)",
    "\u2022  First to identify invisible\n   ground truth as systematic flaw",
    "\u2022  First to map VLM accuracy to\n   Cleveland & McGill hierarchy",
    "\u2022  First design guidelines for\n   VLM + immersive analytics",
    "\u2022  Largest single-benchmark eval\n   (5,713 images with GPT-5.4)",
], font_size=13)


# ==========================================================================
# SLIDE 14: Limitations & Future Work
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Limitations & Future Work", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
             "Limitations", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), [
    "\u2022  Single model (GPT-5.4 only)",
    "\u2022  One question per chart in ChartX",
    "\u2022  3D evaluation limited to bar charts only",
    "\u2022  Automated scoring may miss correct answers",
    "\u2022  Static images, not live VR screenshots",
], font_size=14)

add_text_box(sl, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
             "Future Work", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(sl, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), [
    "\u2022  DXR + Unity immersive rendering\n   condition (Meta Quest)",
    "\u2022  Multi-viewpoint evaluation\n   (3\u20136 camera angles per chart)",
    "\u2022  Structured grounding\n   (image + JSON system state)",
    "\u2022  Multi-model comparison\n   (Claude, Gemini, open-source)",
    "\u2022  Human study: VLM-assisted vs\n   unassisted chart reading in VR",
], font_size=14)


# ==========================================================================
# SLIDE 15: Summary & Thank You
# ==========================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

add_text_box(sl, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Summary", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Three Key Findings", font_size=20, color=ACCENT_BLUE, bold=True)

findings = [
    ("85.7% on 2D, 59.3% on 3D",
     "VLMs are strong on flat charts but collapse under 3D perspective."),
    ("43pp gap follows perceptual science",
     "Position-encoded charts (92\u201395%) >> Area-encoded charts (52\u201371%),\n"
     "exactly as Cleveland & McGill predicted in 1984."),
    ("Invisible ground truth",
     "Treemap failures are partly a benchmark design problem \u2014\n"
     "questions expect data not visible in the image."),
]
for i, (title, desc) in enumerate(findings):
    y = Inches(1.9 + i * 1.2)
    num = str(i + 1)
    shape = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(sl, Inches(1.8), y, Inches(10), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(sl, Inches(1.8), y + Inches(0.4), Inches(10), Inches(0.6),
                 desc, font_size=13, color=LIGHT_GRAY)

# Takeaway
add_text_box(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
             "VLM chart comprehension follows the same perceptual hierarchy as human vision.\n"
             "Position is easy. Area is hard. 3D makes everything worse.",
             font_size=18, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(6.6), Inches(11), Inches(0.6),
             "Thank you!  |  Questions?",
             font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(7), Inches(11), Inches(0.4),
             "nchennu@gmu.edu  |  hbuchire@gmu.edu",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SAVE
# ==========================================================================
output_path = "final_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"  {len(prs.slides)} slides")
