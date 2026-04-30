"""Add presenter notes with equal time split between Naga and Hemanjali."""

from pptx import Presentation
from pptx.util import Pt

prs = Presentation("final_presentation.pptx")

notes = {
    # =====================================================================
    # NAGA's SECTION (Slides 1-7) — ~7.5 minutes
    # =====================================================================

    0: """[NAGA]

Good morning everyone. My name is Naga Venkata Sai Chennu, and this is my partner Hemanjali Buchireddy. Today we are presenting our final project for CS 692, Mobile Immersive Computing, under Dr. Bo Han.

Our project is titled "An Empirical Evaluation of Vision-Language Model Visualization Literacy Across 18 Chart Types." We worked on this in collaboration with Fahim Arsad Nafis, a PhD candidate in Dr. Han's research group.

I will cover the first half of our presentation — the project overview, dataset, methodology, and our main accuracy results. Then Hemanjali will take over for the second half — the 3D results, failure analysis, design guidelines, and future work.

Let's get started.""",

    1: """[NAGA]

Here is a quick overview of our project.

The core research question is: How well can Vision-Language Models understand data visualizations, and what happens when those charts move from 2D to 3D?

This matters because immersive analytics — VR and AR — renders charts in 3D by default. If VLMs cannot read 3D charts, they cannot serve as assistants in VR.

We evaluated GPT-5.4 on the ChartX benchmark — 6,000 chart images across 18 chart types and 22 academic topics. Every image comes with a human-authored question and ground-truth answer.

Looking at the key numbers on the right: 85.7 percent on 2D charts, but only 59.3 percent on 3D. A 43.7 percentage-point gap between the best and worst chart types. And we tested all 6,000 images.

No prior work has evaluated VLMs across this many chart types in a single study.""",

    2: """[NAGA]

Let me tell you about the ChartX dataset.

It is publicly available on HuggingFace under InternScience/ChartX. It contains 6,000 PNG chart images organized into 18 chart types. The topics span 22 academic domains including healthcare, economics, education, and environmental science.

Each image has exactly one question and one answer, authored by human annotators. On the right you can see a sample bar chart. The question is: "How many more nurses are in the East region compared to the North region?" The answer is 500.

These are the original images from the benchmark — we did not generate or modify any of them. This makes our results directly comparable to other researchers using ChartX.""",

    3: """[NAGA]

Our pipeline is fully automated with five steps.

Step 1: Load each chart image and its question from ChartX.
Step 2: Send the image and question to GPT-5.4 with temperature zero for deterministic results.
Step 3: Parse the model's free-text response.
Step 4: Score using 10 percent numeric tolerance and keyword matching.
Step 5: Cache every response as JSON for reproducibility.

This is a zero-shot evaluation — no system prompt, no few-shot examples, no chain-of-thought. We are testing the model's raw chart reading ability. The total cost was approximately 10 dollars for all 6,000 images.""",

    4: """[NAGA]

Now the main results. This is our most important slide.

GPT-5.4 achieves 85.7 percent overall accuracy across 5,713 two-dimensional chart images. But look at the variation.

At the top: line charts at 95.3 percent, pie charts at 94.7 percent, bar charts with labels at 92.1 percent. Thirteen chart types exceed 83 percent.

At the bottom: area charts at 77.4 percent, bubble charts at 71.1 percent, radar at 68.2 percent, and treemaps at just 51.6 percent.

The gap between best and worst is 43.7 percentage points. This is not random — it follows a pattern that perceptual science predicted 40 years ago. Let me show you.""",

    5: """[NAGA]

Why do some chart types score 95 percent and others only 52 percent?

In 1984, Cleveland and McGill ranked how accurately humans read different visual encodings. Position along a common scale is the most accurate. Area is among the least accurate.

Our results match this hierarchy almost perfectly. Position-encoded charts like bar and line average 90 percent. Angle and arc charts like pie average 92.5 percent. But area-encoded charts like treemaps and bubble charts average only about 60 percent.

Stevens' Power Law explains why: the perception exponent for line length is 1.0, meaning perception is linear. For area, the exponent is only 0.7, meaning a rectangle twice the area appears only about 60 percent larger. This makes exact value extraction from area encoding fundamentally imprecise.

VLMs have learned the same perceptual hierarchy as human vision.

Now I will hand over to Hemanjali, who will present the 3D results, failure analysis, and our design guidelines.""",

    # =====================================================================
    # HEMANJALI's SECTION (Slides 7-14) — ~7.5 minutes
    # =====================================================================

    6: """[HEMANJALI]

Thank you, Naga. I will now present our 3D evaluation results.

We evaluated all 280 three-dimensional bar chart images in ChartX. The results are dramatic.

Two-dimensional bar charts with numeric labels achieve 92.1 percent. Without labels, 87.7 percent. But 3D bar charts drop to just 59.3 percent.

That is a 26 to 33 percentage-point collapse, just from changing the rendering from 2D to 3D. The data is the same. The questions are the same. Only the visual presentation changed.

You can see the two sample images here. The 2D chart on the left is clean and easy to read. The 3D chart on the right shows the same data but with perspective distortion and depth that make it much harder to determine exact values.

This has major implications for immersive analytics, where 3D is the default.""",

    7: """[HEMANJALI]

Why exactly do 3D charts cause this failure? Claus Wilke identifies four distortion mechanisms.

First, non-invertible projection. When a 3D chart becomes a 2D image, you lose depth information. You cannot recover exact values from one view.

Second, foreshortening. Bars farther from the camera appear shorter than equally-valued bars in front.

Third, occlusion. Front elements physically hide rear elements. Information is literally removed from the image.

Fourth, unequal scaling. Perspective makes identical data increments occupy different pixel heights at different depths.

The critical point is on the right: VLMs see only a single static 2D projection. Unlike humans in VR, they cannot rotate the view, look behind bars, or use stereo depth. They are stuck with one perspective.""",

    8: """[HEMANJALI]

Now our most important finding — what we call the invisible ground truth problem.

Look at this treemap on the right. What do you see? Colored rectangles and category names. No numbers anywhere.

But the ground truth expects exact percentages — like "20 percent" or "34 percent." These numbers come from the CSV data file used to generate the chart, not from anything visible in the image.

The model is being asked to read numbers that do not exist in the image. Even humans cannot extract exact percentages from rectangle areas without labels. Stevens' Power Law tells us this is a fundamental perceptual limitation, not a model limitation.

ChartX's own paper acknowledges this. They write: "the column label of values is usually invisible" for treemaps. We estimate that correcting for these unanswerable questions would raise treemap accuracy from 52 percent to about 70 to 75 percent.

This is a benchmark design issue, not a VLM failure.""",

    9: """[HEMANJALI]

Let me walk you through the failure analysis for the four weakest chart types.

We classified every incorrect answer into three categories: completely wrong, close but outside tolerance, and format mismatch.

For treemaps, 79 percent of failures are completely wrong — the model is guessing because there are no numbers to read. For radar charts, 22 percent are format mismatches — the model knows the answer but phrases it differently. For bubble charts, the model cannot read values from bubble size. And for area charts, it confuses individual values with cumulative stacked values.

These different failure patterns mean different chart types fail for different reasons. A one-size-fits-all solution will not work.""",

    10: """[HEMANJALI]

Based on all our findings, we propose six design guidelines.

G1: Use position-encoded charts like bar and line for VLM tasks — they achieve 90 percent plus.

G2: Always include numeric labels. Adding labels increases bar chart accuracy by 4.4 percentage points.

G3: Flatten 3D charts before VLM processing. Render 2D projections or provide data tables.

G4: Use tiered scoring tolerance — 5 percent for bar and line, 25 percent for treemap.

G5: Structure prompts as describe, extract, then reason pipelines.

G6: In VR systems, pass DXR JSON specifications alongside viewport images so the VLM has access to the data values and encodings, not just the visual.""",

    11: """[HEMANJALI]

This slide connects our work to immersive analytics and the DXR toolkit.

DXR is a Unity toolkit for building VR data visualizations. You define charts using JSON specifications — chart type, data encodings, interactions — and DXR renders them in 3D for Meta Quest headsets.

Our results predict what happens when a VLM tries to interpret these 3D charts. Based on our findings, a VLM reading a DXR viewport screenshot would achieve about 59 percent on bar charts.

But Guideline 6 offers a solution: the same DXR JSON that defines the chart for rendering can be passed to the VLM as structured context alongside the image. This gives the model both the visual and the data, compensating for information lost in 3D projection.

The architecture is on the right: user sees the chart in VR, the VLM receives the viewport plus DXR JSON, answers the question, and the answer appears as a VR overlay.""",

    12: """[HEMANJALI]

Here is how our work compares to prior studies.

Kim et al. tested GPT-4V on custom tasks and got 60 to 78 percent. Pandey and Ottley tested four VLMs on VLAT and found the same pattern we see — line charts score high, bubble charts score low. Islam et al. did the first multi-benchmark evaluation.

Our contributions are: the broadest chart type diversity at 18 types, the first identification of invisible ground truth as a systematic flaw, the first mapping of VLM accuracy to Cleveland and McGill's hierarchy, the first design guidelines for VLM plus immersive analytics, and the largest single-benchmark evaluation with 5,713 images.""",

    13: """[HEMANJALI]

We want to be transparent about limitations. We tested only GPT-5.4 — cross-model comparison would strengthen the findings. ChartX has only one question per chart. The 3D evaluation covers only bar charts. And we tested static images, not live VR screenshots.

For future work, we plan to use DXR in Unity to render immersive 3D charts for Meta Quest. We want to test multi-viewpoint input with 3 to 6 camera angles. We plan structured grounding with system state JSON. We want multi-model comparison. And eventually a human study comparing VLM-assisted versus unassisted chart reading in VR.

Now I will hand back to Naga for the summary.""",

    14: """[NAGA]

Let me summarize our three key findings.

One: GPT-5.4 achieves 85.7 percent on 2D but drops to 59.3 percent on 3D. VLMs are strong on flat charts but collapse under 3D perspective.

Two: The 43-point gap between chart types follows the perceptual hierarchy. Position-encoded charts score 92 to 95 percent. Area-encoded charts score 52 to 71 percent. Exactly as Cleveland and McGill predicted in 1984.

Three: Some treemap failures are a benchmark problem, not a model problem. The image contains no numbers, but the ground truth expects exact percentages.

Our takeaway: VLM chart comprehension follows the same perceptual hierarchy as human vision. Position is easy. Area is hard. 3D makes everything worse.

Thank you very much. We are happy to take questions.""",
}

for i, slide in enumerate(prs.slides):
    if i in notes:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes[i]
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14)

prs.save("final_presentation.pptx")
print("Speaker notes updated with equal time split:")
print("  Naga:     Slides 1-6, 15 (~7.5 min)")
print("  Hemanjali: Slides 7-14 (~7.5 min)")
print("Saved: final_presentation.pptx")
